# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建PPT对象
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色主题
PRIMARY_COLOR = RGBColor(0x2E, 0x86, 0xAB)  # 蓝色
SECONDARY_COLOR = RGBColor(0x06, 0xD6, 0xA0)  # 绿色
ACCENT_COLOR = RGBColor(0xFF, 0xD1, 0x66)  # 黄色
DARK_COLOR = RGBColor(0x2D, 0x3A, 0x4A)  # 深蓝
LIGHT_COLOR = RGBColor(0xF8, 0xF9, 0xFA)  # 浅灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PYTHON_BLUE = RGBColor(0x37, 0x76, 0xAB)
PYTHON_YELLOW = RGBColor(0xFF, 0xD4, 0x3B)

def add_background(slide, color):
    """设置幻灯片背景颜色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None):
    """添加形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, font_color=DARK_COLOR, bold=False, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox

def add_code_block(slide, left, top, width, height, code, font_size=14):
    """添加代码块"""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color=RGBColor(0x28, 0x2C, 0x34))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xAB, 0xB2, 0xBF)
    p.font.name = "Consolas"
    return shape

def create_title_slide():
    """创建封面幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局

    # 添加装饰形状
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill_color=DARK_COLOR)

    # 添加Python logo装饰
    add_shape(slide, MSO_SHAPE.OVAL, Inches(9.5), Inches(0.5), Inches(3), Inches(3), fill_color=PYTHON_YELLOW)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(9.8), Inches(0.8), Inches(2.4), Inches(2.4), fill_color=PYTHON_BLUE)

    # 添加装饰线条
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1), Inches(3.2), Inches(8), Inches(0.05), fill_color=SECONDARY_COLOR)

    # 标题
    add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(1.5),
                "Python 编程学习", font_size=54, font_color=WHITE, bold=True)

    # 副标题
    add_textbox(slide, Inches(1), Inches(3.5), Inches(8), Inches(1),
                "从零开始掌握Python编程语言", font_size=28, font_color=ACCENT_COLOR)

    # 版本信息
    add_textbox(slide, Inches(1), Inches(5), Inches(8), Inches(0.8),
                "Python 3.x | 基础语法 | 数据结构 | 函数模块", font_size=18, font_color=LIGHT_COLOR)

    # 底部信息
    add_textbox(slide, Inches(1), Inches(6.5), Inches(8), Inches(0.5),
                "Duron204 制作", font_size=16, font_color=RGBColor(0x6C, 0x75, 0x7D))

def create_toc_slide():
    """创建目录幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # 标题
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.2), fill_color=PRIMARY_COLOR)
    add_textbox(slide, Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
                "目录 CONTENTS", font_size=36, font_color=WHITE, bold=True)

    # 目录内容
    toc_items = [
        ("01", "Python 基础入门", "第一个程序 | 变量与数据类型"),
        ("02", "运算符详解", "算术 | 赋值 | 比较 | 逻辑运算符"),
        ("03", "流程控制", "分支结构 | 循环结构 | 实战练习"),
        ("04", "数据结构", "列表 | 元组 | 字符串 | 集合 | 字典"),
        ("05", "函数与模块", "定义函数 | 参数 | 模块管理"),
    ]

    for i, (num, title, desc) in enumerate(toc_items):
        y_pos = Inches(1.5 + i * 1.1)

        # 数字圆形
        circle = add_shape(slide, MSO_SHAPE.OVAL, Inches(1), y_pos, Inches(0.7), Inches(0.7), fill_color=SECONDARY_COLOR)
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 标题
        add_textbox(slide, Inches(2), y_pos, Inches(4), Inches(0.5),
                    title, font_size=22, font_color=DARK_COLOR, bold=True)

        # 描述
        add_textbox(slide, Inches(6), y_pos, Inches(6), Inches(0.5),
                    desc, font_size=16, font_color=RGBColor(0x6C, 0x75, 0x7D))

def create_section_slide(section_num, title, subtitle):
    """创建章节封面幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill_color=PRIMARY_COLOR)

    # 装饰圆形
    add_shape(slide, MSO_SHAPE.OVAL, Inches(8), Inches(-1), Inches(6), Inches(6), fill_color=RGBColor(0x3A, 0x9B, 0xBD))
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(4), Inches(4), fill_color=RGBColor(0x20, 0x7A, 0x9E))

    # 章节编号
    add_textbox(slide, Inches(1), Inches(1), Inches(4), Inches(1.5),
                section_num, font_size=72, font_color=ACCENT_COLOR, bold=True)

    # 标题
    add_textbox(slide, Inches(1), Inches(3), Inches(10), Inches(1.5),
                title, font_size=48, font_color=WHITE, bold=True)

    # 副标题
    add_textbox(slide, Inches(1), Inches(5), Inches(10), Inches(1),
                subtitle, font_size=24, font_color=RGBColor(0xCC, 0xE5, 0xFF))

def create_content_slide(title, content_pairs, has_code=True):
    """创建内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # 顶部标题栏
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1), fill_color=PRIMARY_COLOR)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                title, font_size=28, font_color=WHITE, bold=True)

    # 左侧内容区域
    left_width = Inches(6) if has_code else Inches(12)

    # 添加内容
    y_pos = Inches(1.3)
    for i, (text, code) in enumerate(content_pairs):
        # 文本说明
        add_textbox(slide, Inches(0.5), y_pos, left_width, Inches(0.8),
                    text, font_size=16, font_color=DARK_COLOR)

        if has_code and code:
            # 代码块
            add_code_block(slide, Inches(6.8), y_pos, Inches(6), Inches(min(1.5, 0.4 * len(code.split('\n')))), code, font_size=12)

        y_pos += Inches(1.2)

    return slide

def create_code_slide(title, code_blocks):
    """创建代码展示幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # 顶部标题栏
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1), fill_color=DARK_COLOR)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                title, font_size=28, font_color=PYTHON_YELLOW, bold=True)

    # Python标志
    add_textbox(slide, Inches(12), Inches(0.15), Inches(1), Inches(0.7),
                ">>>", font_size=24, font_color=SECONDARY_COLOR, bold=True, alignment=PP_ALIGN.RIGHT)

    # 代码块
    y_pos = Inches(1.3)
    for desc, code in code_blocks:
        if desc:
            add_textbox(slide, Inches(0.5), y_pos, Inches(12), Inches(0.5),
                        desc, font_size=16, font_color=PRIMARY_COLOR, bold=True)
            y_pos += Inches(0.5)

        code_lines = len(code.split('\n'))
        code_height = min(Inches(3), Inches(0.3 * code_lines))
        add_code_block(slide, Inches(0.5), y_pos, Inches(12.3), code_height, code, font_size=13)
        y_pos += code_height + Inches(0.3)

    return slide

def create_tip_slide(tips):
    """创建提示/技巧幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, LIGHT_COLOR)

    # 标题
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1), fill_color=SECONDARY_COLOR)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                "学习要点与技巧", font_size=32, font_color=WHITE, bold=True)

    # 添加提示卡片
    for i, (icon, title, desc) in enumerate(tips):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.5)
        y = Inches(1.3 + row * 2.5)

        # 卡片背景
        card = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(2.2), fill_color=WHITE)

        # 图标
        icon_shape = add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.3), y + Inches(0.3), Inches(0.8), Inches(0.8), fill_color=ACCENT_COLOR)
        tf = icon_shape.text_frame
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

        # 标题
        add_textbox(slide, x + Inches(1.3), y + Inches(0.3), Inches(4.5), Inches(0.5),
                    title, font_size=18, font_color=DARK_COLOR, bold=True)

        # 描述
        add_textbox(slide, x + Inches(1.3), y + Inches(0.9), Inches(4.5), Inches(1),
                    desc, font_size=14, font_color=RGBColor(0x6C, 0x75, 0x7D))

    return slide

def create_practice_slide(title, exercises):
    """创建练习题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, WHITE)

    # 标题
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1), fill_color=ACCENT_COLOR)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.7),
                title, font_size=28, font_color=DARK_COLOR, bold=True)

    # 练习题
    for i, (num, question, hint) in enumerate(exercises):
        y = Inches(1.3 + i * 1.5)

        # 编号
        num_shape = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.6), Inches(0.6), fill_color=PRIMARY_COLOR)
        tf = num_shape.text_frame
        p = tf.paragraphs[0]
        p.text = str(num)
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 问题
        add_textbox(slide, Inches(1.3), y, Inches(8), Inches(0.5),
                    question, font_size=16, font_color=DARK_COLOR, bold=True)

        # 提示
        if hint:
            add_textbox(slide, Inches(1.3), y + Inches(0.5), Inches(10), Inches(0.8),
                        f"提示: {hint}", font_size=13, font_color=RGBColor(0x6C, 0x75, 0x7D))

    return slide

def create_summary_slide(points):
    """创建总结幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_COLOR)

    # 装饰
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10), Inches(-2), Inches(5), Inches(5), fill_color=RGBColor(0x3A, 0x4A, 0x5C))

    # 标题
    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(1),
                "本章小结", font_size=40, font_color=ACCENT_COLOR, bold=True)

    # 分割线
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(3), Inches(0.05), fill_color=SECONDARY_COLOR)

    # 要点
    for i, point in enumerate(points):
        y = Inches(2 + i * 0.8)

        # 检查标记
        check = add_shape(slide, MSO_SHAPE.OVAL, Inches(0.5), y, Inches(0.4), Inches(0.4), fill_color=SECONDARY_COLOR)
        tf = check.text_frame
        p = tf.paragraphs[0]
        p.text = "✓"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        add_textbox(slide, Inches(1.1), y, Inches(11), Inches(0.6),
                    point, font_size=18, font_color=WHITE)

    return slide

def create_end_slide():
    """创建结束幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, PRIMARY_COLOR)

    # 装饰
    add_shape(slide, MSO_SHAPE.OVAL, Inches(5), Inches(1), Inches(3.333), Inches(3.333), fill_color=PYTHON_YELLOW)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(5.3), Inches(1.3), Inches(2.733), Inches(2.733), fill_color=PYTHON_BLUE)

    # 主标题
    add_textbox(slide, Inches(0), Inches(4.5), Inches(13.333), Inches(1),
                "Thank You", font_size=56, font_color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 副标题
    add_textbox(slide, Inches(0), Inches(5.5), Inches(13.333), Inches(0.8),
                "继续学习，不断进步！", font_size=24, font_color=ACCENT_COLOR, alignment=PP_ALIGN.CENTER)

    # 底部信息
    add_textbox(slide, Inches(0), Inches(6.5), Inches(13.333), Inches(0.5),
                "Python 编程学习 | Duron204", font_size=16, font_color=RGBColor(0xCC, 0xE5, 0xFF), alignment=PP_ALIGN.CENTER)

# ==================== 开始创建PPT ====================

# 1. 封面
create_title_slide()

# 2. 目录
create_toc_slide()

# 3. 第一章：Python基础入门
create_section_slide("01", "Python 基础入门", "从第一个程序开始")

create_code_slide("第一个Python程序", [
    ("输出函数 print()", 'print("Hello, World!")\nprint("你好，世界！")'),
    ("运行结果", 'Hello, World!\n你好，世界！'),
])

create_content_slide("变量与数据类型", [
    ("整数类型 (int)", "a = 100"),
    ("浮点数类型 (float)", "b = 3.14"),
    ("字符串类型 (str)", 'c = "Hello"'),
    ("布尔类型 (bool)", "d = True / False"),
    ("查看变量类型", "print(type(a))  # <class 'int'>"),
], has_code=False)

create_code_slide("变量与类型转换", [
    ("定义变量", 'a = 100\nb = 123.45\nc = "hello"\nd = True'),
    ("类型查看", "print(type(a))  # <class 'int'>\nprint(type(b))  # <class 'float'>"),
    ("类型转换", 'print(float(a))    # 100.0\nprint(int(b))      # 123\nprint(str(a))      # "100"'),
])

create_tip_slide([
    ("1", "变量命名规则", "以字母或下划线开头，区分大小写，避免使用关键字"),
    ("2", "数据类型", "Python是动态类型语言，变量类型由赋值决定"),
    ("3", "类型转换", "使用int()、float()、str()等函数进行类型转换"),
    ("4", "注释", "单行注释用#，多行注释用三引号"),
])

# 4. 第二章：运算符
create_section_slide("02", "运算符详解", "掌握Python中的各类运算符")

create_code_slide("算术运算符", [
    ("基本运算", 'print(321 + 12)   # 加法: 333\nprint(321 - 12)   # 减法: 309\nprint(321 * 12)   # 乘法: 3852'),
    ("除法和取整", 'print(321 / 12)   # 除法: 26.75\nprint(321 // 12)  # 整除: 26\nprint(321 % 12)   # 取余: 9'),
    ("幂运算", 'print(321 ** 2)   # 幂运算: 103041'),
])

create_code_slide("赋值与比较运算符", [
    ("赋值运算符", 'a = 10\na += 5   # a = a + 5\na *= 2   # a = a * 2'),
    ("比较运算符", 'print(3 > 2)    # True\nprint(3 == 3)   # True\nprint(3 != 2)   # True'),
    ("逻辑运算符", 'print(True and False)  # False\nprint(True or False)   # True\nprint(not True)        # False'),
])

create_practice_slide("运算符练习", [
    (1, "计算圆的面积和周长", "输入半径，使用math.pi计算"),
    (2, "华氏温度转摄氏温度", "公式: C = (F - 32) / 1.8"),
    (3, "判断闰年", "能被4整除但不能被100整除，或能被400整除"),
])

# 5. 第三章：流程控制
create_section_slide("03", "流程控制", "分支与循环结构")

create_code_slide("if 分支结构", [
    ("简单条件判断", 'age = 18\nif age >= 18:\n    print("已成年")\nelse:\n    print("未成年")'),
    ("多条件判断", 'score = 85\nif score >= 90:\n    grade = "A"\nelif score >= 80:\n    grade = "B"\nelif score >= 70:\n    grade = "C"\nelse:\n    grade = "D"'),
])

create_code_slide("for 循环结构", [
    ("基本for循环", 'for i in range(5):\n    print(i)  # 0, 1, 2, 3, 4'),
    ("遍历列表", 'fruits = ["apple", "banana", "cherry"]\nfor fruit in fruits:\n    print(fruit)'),
    ("列表推导式", 'squares = [x**2 for x in range(10)]\nprint(squares)  # [0, 1, 4, 9, 16, 25, 36, 49, 64, 81]'),
])

create_code_slide("while 循环与控制", [
    ("while循环", 'total = 0\ni = 1\nwhile i <= 100:\n    total += i\n    i += 1\nprint(total)  # 5050'),
    ("break和continue", 'for i in range(10):\n    if i == 3:\n        continue  # 跳过3\n    if i == 7:\n        break     # 到7停止\n    print(i)'),
])

create_practice_slide("流程控制练习", [
    (1, "BMI计算器", "根据身高体重计算BMI并判断体型"),
    (2, "打印九九乘法表", "使用嵌套循环"),
    (3, "猜数字游戏", "随机生成数字，用户猜测"),
    (4, "判断素数", "判断输入的数是否为素数"),
])

# 6. 第四章：数据结构
create_section_slide("04", "数据结构", "列表 | 元组 | 字符串 | 集合 | 字典")

create_code_slide("列表 List", [
    ("创建列表", 'fruits = ["apple", "banana", "cherry"]\nnumbers = [1, 2, 3, 4, 5]'),
    ("列表操作", 'fruits.append("orange")   # 添加元素\nfruits.remove("banana")  # 删除元素\nprint(len(fruits))       # 获取长度'),
    ("列表切片", 'nums = [0, 1, 2, 3, 4, 5]\nprint(nums[1:4])    # [1, 2, 3]\nprint(nums[::-1])   # [5, 4, 3, 2, 1, 0]'),
])

create_code_slide("元组 Tuple", [
    ("创建元组", 'point = (3, 4)\ncolors = ("red", "green", "blue")'),
    ("元组特性", '# 元组是不可变的\n# point[0] = 5  # TypeError\n\n# 解包\nx, y = point\nprint(x, y)  # 3 4'),
])

create_code_slide("字符串 String", [
    ("字符串方法", 's = "Hello, World!"\nprint(s.upper())      # HELLO, WORLD!\nprint(s.lower())      # hello, world!\nprint(s.replace("World", "Python"))'),
    ("字符串格式化", 'name = "Python"\nversion = 3\nprint(f"{name} {version}")  # Python 3\nprint("{} {}".format(name, version))'),
])

create_code_slide("集合 Set 与字典 Dict", [
    ("集合操作", 's1 = {1, 2, 3}\ns2 = {3, 4, 5}\nprint(s1 | s2)  # 并集: {1, 2, 3, 4, 5}\nprint(s1 & s2)  # 交集: {3}'),
    ("字典操作", 'person = {"name": "Alice", "age": 25}\nprint(person["name"])      # Alice\nperson["email"] = "a@b.com"  # 添加\nfor k, v in person.items():\n    print(f"{k}: {v}")'),
])

create_tip_slide([
    ("1", "列表 vs 元组", "列表可变，元组不可变；元组性能更好"),
    ("2", "字符串方法", "字符串是不可变的，方法返回新字符串"),
    ("3", "集合特性", "集合元素唯一，支持集合运算"),
    ("4", "字典键", "字典键必须是不可变类型"),
])

# 7. 第五章：函数与模块
create_section_slide("05", "函数与模块", "代码复用与模块化编程")

create_code_slide("定义和调用函数", [
    ("基本函数", 'def greet(name):\n    return f"Hello, {name}!"\n\nprint(greet("Python"))  # Hello, Python!'),
    ("默认参数", 'def power(base, exp=2):\n    return base ** exp\n\nprint(power(3))      # 9\nprint(power(3, 3))   # 27'),
])

create_code_slide("参数类型", [
    ("位置参数", 'def add(a, b):\n    return a + b\n\nprint(add(1, 2))  # 3'),
    ("可变参数", 'def sum_all(*args):\n    return sum(args)\n\nprint(sum_all(1, 2, 3, 4))  # 10'),
    ("关键字参数", 'def info(**kwargs):\n    for k, v in kwargs.items():\n        print(f"{k}: {v}")\n\ninfo(name="Alice", age=25)'),
])

create_code_slide("模块管理", [
    ("导入模块", 'import math\nprint(math.pi)  # 3.14159...'),
    ("导入特定函数", 'from math import sqrt, pi\nprint(sqrt(16))  # 4.0'),
    ("别名导入", 'import numpy as np\nfrom math import factorial as fac'),
])

create_practice_slide("函数练习", [
    (1, "编写阶乘函数", "使用循环或递归实现"),
    (2, "编写斐波那契函数", "返回第n个斐波那契数"),
    (3, "编写装饰器", "实现函数执行计时"),
])

# 8. 总结
create_summary_slide([
    "掌握了Python基础语法和变量类型",
    "学会了使用运算符进行各种计算",
    "能够使用分支和循环控制程序流程",
    "理解了常用数据结构的特性和用法",
    "掌握了函数定义和模块化编程",
    "通过实战练习巩固了所学知识",
])

# 9. 结束页
create_end_slide()

# 保存PPT
prs.save("d:/VScode project/python学习/Python学习指南.pptx")
print("PPT创建成功！")
